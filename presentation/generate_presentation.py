"""Generate KSAS Spring 2026 conference PowerPoint from the blue Canva template."""

import os
import subprocess
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(BASE, "fig")
TEMPLATE = os.path.join(BASE, "template.pptx")

# Use a Korean-capable bold font for all generated text.
FONT_NAME = "NanumSquare Bold"

# ── Colors (from template: grey bg, white panel, blue+black text) ───

BG_GREY = RGBColor(0xDC, 0xDC, 0xDC)
BLUE = RGBColor(0x00, 0x4A, 0xAD)
BLUE_PALE = RGBColor(0xD0, 0xE0, 0xF5)
HEADER_PALE = RGBColor(0xEA, 0xF3, 0xFF)  # template header block color
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GREY = RGBColor(0x66, 0x66, 0x66)
KO_GREY = RGBColor(0x99, 0x99, 0x99)
LIGHT = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
RED_LIGHT = RGBColor(0xFA, 0xE5, 0xE3)

# ── Dimensions (template 20" × 11.25") ─────────────────────────────

W = Inches(20)
H = Inches(11.25)
# White content panel matching template's Group 2 (full-height variant)
PL, PT_, PW, PH = Inches(0.88), Inches(0.83), Inches(18.23), Inches(9.44)
# Header band occupies top 2.11" of the white panel
HEADER_H = Inches(2.11)
# Content zone starts right below the header
CONTENT_TOP = Inches(0.83 + 2.11 + 0.3)  # ~3.24"
CL = Inches(1.6)
CW = Inches(16.8)
TOTAL = 17


# ── Helpers ─────────────────────────────────────────────────────────

def _rr(sl, l, t, w, h, fill, line=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _rounded_rect(sl, l, t, w, h, fill, line=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _tb(sl, l, t, w, h):
    tb = sl.shapes.add_textbox(l, t, w, h)
    # Set default font on the first paragraph so direct `p = ...paragraphs[0]`
    # formatting later inherits the same font.
    try:
        p0 = tb.text_frame.paragraphs[0]
        p0.font.name = FONT_NAME
        p0.font.bold = True
    except Exception:
        pass
    return tb


def _set(shape, text, sz=22, color=DARK, bold=True, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.name = FONT_NAME
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def _p(tf, text, sz=22, color=DARK, bold=True, sp=Pt(8), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.name = FONT_NAME
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = sp
    p.alignment = align
    return p


def _ko(tf, text, sz=20, sp=Pt(2), align=PP_ALIGN.LEFT):
    return _p(tf, text, sz, KO_GREY, True, sp, align)


def _bullet(tf, en, kr=None, sz=20, sp=Pt(10)):
    _p(tf, f"\u2022  {en}", sz, DARK, sp=sp)
    if kr:
        _ko(tf, f"    {kr}", sz, Pt(2))


def _img(sl, path, l, t, width=None, height=None):
    if os.path.exists(path):
        try:
            return sl.shapes.add_picture(path, l, t, width, height)
        except ValueError as exc:
            # Some phone-camera JPGs are MPO containers. Convert to PNG fallback.
            if "unsupported image format" in str(exc):
                try:
                    from PIL import Image
                    with Image.open(path) as im:
                        im = im.convert("RGB")
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                            tmp_path = tmp.name
                        im.save(tmp_path, format="PNG")
                    try:
                        return sl.shapes.add_picture(tmp_path, l, t, width, height)
                    finally:
                        if os.path.exists(tmp_path):
                            os.unlink(tmp_path)
                except Exception as conv_exc:
                    print("image conversion failed", path, conv_exc)
            print("image insert failed", path, exc)
    if not os.path.exists(path):
        print("path not found", path)
    else:
        print("using placeholder for image", path)
    box = _rr(sl, l, t, width or Inches(6), height or Inches(4), LIGHT, GREY)
    _set(box, f"[{os.path.basename(path)}]", 20, GREY, align=PP_ALIGN.CENTER)
    return box


def _video(sl, video_path, l, t, width, height, poster_path=None):
    if not os.path.exists(video_path):
        print("video path not found", video_path)
        box = _rr(sl, l, t, width, height, LIGHT, GREY)
        _set(box, f"[{os.path.basename(video_path)}]", 20, GREY, align=PP_ALIGN.CENTER)
        return box

    # Some python-pptx builds can fail on movie insertion; keep slide generation resilient.
    try:
        if poster_path and os.path.exists(poster_path):
            try:
                return sl.shapes.add_movie(
                    video_path,
                    l,
                    t,
                    width,
                    height,
                    poster_frame_image=poster_path,
                    mime_type="video/mp4",
                )
            except ValueError as exc:
                # Retry without poster when camera JPG is unsupported (e.g., MPO).
                if "unsupported image format" not in str(exc):
                    raise
        return sl.shapes.add_movie(video_path, l, t, width, height, mime_type="video/mp4")
    except Exception as exc:
        print("video embed failed", video_path, exc)
        box = _rr(sl, l, t, width, height, LIGHT, GREY)
        _set(box, "[Insert demo video manually]", 20, GREY, align=PP_ALIGN.CENTER)
        return box


def _export_rotated_png(path, rotate_cw_deg):
    """Write a PNG copy of *path* rotated *rotate_cw_deg* clockwise (0, 90, 180, 270)."""
    from PIL import Image

    with Image.open(path) as im:
        im = im.convert("RGB")
        if rotate_cw_deg == 90:
            im = im.transpose(Image.Transpose.ROTATE_270)
        elif rotate_cw_deg == 180:
            im = im.transpose(Image.Transpose.ROTATE_180)
        elif rotate_cw_deg == 270:
            im = im.transpose(Image.Transpose.ROTATE_90)
        elif rotate_cw_deg != 0:
            raise ValueError(f"unsupported rotate_cw_deg={rotate_cw_deg}")
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        im.save(tmp_path, format="PNG")
    return tmp_path


def _extract_video_poster_frame(video_path, timestamp_s=1.0):
    """Extract a PNG poster frame from a video using ffmpeg."""
    try:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        cmd = [
            "ffmpeg",
            "-y",
            "-ss",
            str(timestamp_s),
            "-i",
            video_path,
            "-frames:v",
            "1",
            "-q:v",
            "2",
            tmp_path,
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if os.path.exists(tmp_path) and os.path.getsize(tmp_path) > 0:
            return tmp_path
    except Exception as exc:
        print("video poster extraction failed", video_path, exc)
    if "tmp_path" in locals() and os.path.exists(tmp_path):
        os.unlink(tmp_path)
    return None


def _add_picture_fit_center(sl, path, l, t, box_w, box_h):
    """Place image in box (l,t,box_w,box_h), preserve aspect ratio, center."""
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    l_in = _emu_to_in(int(l))
    t_in = _emu_to_in(int(t))
    bw_in = _emu_to_in(int(box_w))
    bh_in = _emu_to_in(int(box_h))
    s = min(bw_in / iw, bh_in / ih)
    w_in = iw * s
    h_in = ih * s
    left_in = l_in + (bw_in - w_in) / 2
    top_in = t_in + (bh_in - h_in) / 2
    return sl.shapes.add_picture(
        path, Inches(left_in), Inches(top_in), width=Inches(w_in), height=Inches(h_in)
    )


def _emu_to_in(v):
    return v / 914400.0


def _sn(sl, n):
    tb = _tb(sl, Inches(18.2), Inches(10.55), Inches(1.2), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{n}/{TOTAL}"
    p.font.size = Pt(20)
    p.font.color.rgb = GREY
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT


def _circle(sl, l, t, d, fill, text, text_color=WHITE, text_sz=20):
    c = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    c.shadow.inherit = False
    _set(c, text, text_sz, text_color, True, PP_ALIGN.CENTER)
    c.text_frame.margin_top = Pt(2)
    c.text_frame.margin_left = Pt(0)
    c.text_frame.margin_right = Pt(0)
    return c


def make_slide(prs, layout):
    """Grey bg + white content panel, matching the template."""
    sl = prs.slides.add_slide(layout)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG_GREY
    _rr(sl, PL, PT_, PW, PH, WHITE)
    return sl


F2F2 = RGBColor(0xF2, 0xF2, 0xF2)


def slide_title(sl, text, y=PT_):
    """Template #EAF3FF header band (18.23" x 2.11") with centered title at y=1.82"."""
    header = _rr(sl, PL, y, PW, HEADER_H, HEADER_PALE)
    header.shadow.inherit = False

    tb = _tb(sl, Inches(5.34), Inches(1.5), Inches(9.31), Inches(0.73))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(40)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def _callout(sl, l, t, w, h, text, sz=20):
    """Template-style #EAF3FF rounded callout box with bold text."""
    box = _rounded_rect(sl, l, t, w, h, HEADER_PALE)
    box.shadow.inherit = False
    tb = _tb(sl, l + Inches(0.4), t + Inches(0.4), w - Inches(0.8), h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return tf


def _card(sl, l, t, w, h, fill=F2F2):
    """Template-style #F2F2F2 rounded card background."""
    c = _rounded_rect(sl, l, t, w, h, fill)
    c.shadow.inherit = False
    return c


ARCH_IMG_W_PX = 1646
ARCH_STEPS = [
    (1, 190, 180, "X-Plane"),
    (2, 190, 555, "Host"),
    (3, 730, 850, "Stewart"),
    (4, 1420, 555, "PX4"),
]


def _architecture_figure(sl, l, t, height=Cm(13.0), active_step=None):
    fig1 = os.path.join(FIG, "Fig1_Architecture.png")
    pic = _img(sl, fig1, l, t, height=height)

    img_w_in = _emu_to_in(pic.width) if hasattr(pic, "width") and pic.width else 10.5
    px_per_inch = ARCH_IMG_W_PX / img_w_in
    for num, px_x, px_y, _ in ARCH_STEPS:
        is_active = active_step is None or num == active_step
        cd = Inches(0.75) if is_active else Inches(0.55)
        sz = 24 if is_active else 20
        cx = l + Inches(px_x / px_per_inch) - cd / 2
        cy = t + Inches(px_y / px_per_inch) - cd / 2
        c = _circle(
            sl, cx, cy, cd,
            RED if is_active else BLUE_PALE,
            str(num),
            WHITE if is_active else BLUE,
            sz,
        )
        c.line.color.rgb = WHITE if is_active else BLUE
        c.line.width = Pt(3) if is_active else Pt(1)
    return pic



def _block_slide_base(prs, layout, slide_num, title, step_num):
    """Shared skeleton for block-detail slides. Returns (sl, card_tf).

    Left side: architecture figure with the active step highlighted.
    Right side: card for per-slide content.
    """
    sl = make_slide(prs, layout)
    slide_title(sl, title)

    _architecture_figure(sl, Inches(1.2), Inches(3.4), height=Cm(14.5),
                         active_step=step_num)

    card_l, card_t = Inches(9.8), Inches(3.4)
    card_w, card_h = Inches(8.8), Inches(6.5)
    _card(sl, card_l, card_t, card_w, card_h)

    ctb = _tb(sl, card_l + Inches(0.35), card_t + Inches(0.25),
              card_w - Inches(0.7), card_h - Inches(0.5))
    ctf = ctb.text_frame
    ctf.word_wrap = True

    _sn(sl, slide_num)
    return sl, ctf


# ── Load template, clear slides ─────────────────────────────────────

prs = Presentation(TEMPLATE)
while len(prs.slides._sldIdLst) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

blank = prs.slide_layouts[6]


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)

tb = _tb(sl, Inches(2), Inches(1.2), Inches(16), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.text = "PHYSICALLY CLOSED-LOOP HILS"
p.font.size = Pt(24)
p.font.color.rgb = BLUE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

tb = _tb(sl, Inches(2), Inches(2.0), Inches(16), Inches(3.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Physically Closed-Loop HILS for\nFlight Attitude-Control Validation"
p.font.size = Pt(56)
p.font.color.rgb = BLACK
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(68)

_p(tf, "물리적 폐루프 HILS를 이용한 비행 자세제어 검증", 32, DARK, sp=Pt(18), align=PP_ALIGN.CENTER)
_p(tf, "Estimator-in-the-Loop Validation Using a Stewart Platform, X-Plane, and PX4", 22, GREY, sp=Pt(12), align=PP_ALIGN.CENTER)

tb = _tb(sl, Inches(2), Inches(7.0), Inches(16), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Spring 2026 KSAS Conference"
p.font.size = Pt(22)
p.font.color.rgb = BLACK
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
_p(tf, "\uc815\ud604\uc6a9 \u2014 \uad6d\ubbfc\ub300\ud559\uad50 \ubbf8\ub798\ubaa8\ube4c\ub9ac\ud2f0\uc81c\uc5b4\uc5f0\uad6c\uc2e4", 20, GREY, sp=Pt(6), align=PP_ALIGN.CENTER)

_sn(sl, 1)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — Motivation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "Why Add Physical Motion to HILS?")

ko_tb = _tb(sl, PL, Inches(2.65), PW, Inches(0.5))

col_w = Inches(7.5)
col_top = Inches(3.4)
lx = Inches(1.8)
rx = Inches(10.7)

# Conventional HILS column
ltb = _tb(sl, lx, col_top, col_w, Inches(5.5))
ltf = ltb.text_frame
ltf.word_wrap = True
p = ltf.paragraphs[0]
p.text = "Conventional HILS"
p.font.size = Pt(28)
p.font.color.rgb = BLUE
p.font.bold = True
for en, kr in [
    ("Simulator computes dynamics", "시뮬레이터가 동역학 계산"),
    ("Controller receives synthetic sensor signals", "합성 센서 신호 수신"),
    ("Attitude evaluated via logs / screens", "로그/화면으로만 자세 확인"),
    ("IMU does not experience real motion", "IMU가 실제 관성 운동 미경험"),
]:
    _bullet(ltf, en, kr, 20, Pt(14))

# This Work column
rtb = _tb(sl, rx, col_top, col_w, Inches(5.5))
rtf = rtb.text_frame
rtf.word_wrap = True
p = rtf.paragraphs[0]
p.text = "This Work"
p.font.size = Pt(28)
p.font.color.rgb = BLUE
p.font.bold = True
for en, kr in [
    ("Simulator drives real Stewart platform motion", "실제 플랫폼 모션 구동"),
    ("PX4 IMU senses real inertial excitation", "PX4 IMU 실제 관성 자극 감지"),
    ("Measured attitude fed back into loop", "측정 자세 → 제어 루프 피드백"),
    ("Timing and fidelity measured", "타이밍·충실도 정량 측정"),
]:
    _bullet(rtf, en, kr, 20, Pt(14))

# divider line between columns
div = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.85), Inches(3.6), Pt(2), Inches(4.8))
div.fill.solid()
div.fill.fore_color.rgb = BLUE
div.line.fill.background()

# bottom callout
_callout(sl, Inches(2.67), Inches(8.8), Inches(14.66), Inches(1.2),"물리 폐루프 시험 환경이 안정적으로 동작하고 추종 성능을 정량화할 수 있음을 실험 데이터로 입증", 20)

_sn(sl, 2)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — What Is New
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "What Is New Here?")

ctb = _tb(sl, Inches(2), Inches(3.2), Inches(16), Inches(1.2))
tf = ctb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("PX4 추정기가 합성 신호가 아닌 실제 물리적 관성 운동에 반응 \u2014 검증의 질이 달라진다")
p.font.size = Pt(22)
p.font.color.rgb = DARK
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

bw = Inches(5.0)
bh = Inches(3.2)
btop = Inches(4.8)
gap = Inches(0.5)

total_w = 3 * bw + 2 * gap
left = (prs.slide_width - total_w) // 2
bxs = [left + i * (bw + gap) for i in range(3)]

labels = [
    ("01", "Realized Motion", "모션 구현",
     "Simulator attitude is converted\ninto Stewart-platform motion.",
     "시뮬레이터 자세 → 플랫폼 모션"),
    ("02", "Physical Sensing", "물리적 센싱",
     "PX4 IMU and estimator respond\nto real inertial excitation.",
     "PX4 IMU가 실제 관성 자극에 반응"),
    ("03", "Closed-Loop Return", "폐루프 귀환",
     "Measured attitude is sent back\ninto the host control loop.",
     "측정된 자세 → 호스트 제어 루프"),
]

for i, (num, t_en, t_kr, d_en, d_kr) in enumerate(labels):
    x = bxs[i]

    ntb = _tb(sl, x, btop, bw, Inches(0.5))
    p = ntb.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(30)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    ttb = _tb(sl, x, btop + Inches(0.6), bw, Inches(0.5))
    p = ttb.text_frame.paragraphs[0]
    p.text = t_en
    p.font.size = Pt(22)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    dtb = _tb(sl, x, btop + Inches(1.1), bw, Inches(1.8))
    dtf = dtb.text_frame
    dtf.word_wrap = True
    p = dtf.paragraphs[0]
    p.text = d_en
    p.font.size = Pt(20)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    _ko(dtf, d_kr, 20, Pt(8), PP_ALIGN.CENTER)

ctf = _callout(
    sl, Inches(2.67), Inches(8.5), Inches(14.66), Inches(1.6),
    "스튜어트 플랫폼은 시뮬레이터와 비행 제어기 사이의 물리적 연결 고리", 18
)
_p(ctf, "소프트웨어 자세 명령을 실제 6-DOF 모션으로 변환해 PX4 IMU가 실제 운동을 경험",
   20, GREY, sp=Pt(6), align=PP_ALIGN.CENTER)

_sn(sl, 3)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — System in Action
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "System in Action")

photo_path = os.path.join(FIG, "PX4_on_STP.JPG")
video_path = os.path.join(FIG, "Demo_0324_vid.mp4")

# Lead with the hardware photo on the left, then let the demo video dominate the slide.
photo_card_left = Inches(1.95)
photo_card_top = Inches(3.35)
photo_card_w = Inches(3.35)
photo_card_h = Inches(5.45)

video_left = Inches(5.65)
video_top = photo_card_top
video_w = Inches(12.15)
video_h = photo_card_h

photo_png = None
poster_png = None
try:
    poster_png = _extract_video_poster_frame(video_path, timestamp_s=1.0)
    photo_png = _export_rotated_png(photo_path, 90)

    ptb = _tb(sl, photo_card_left, photo_card_top - Inches(0.45), photo_card_w, Inches(0.35))
    p = ptb.text_frame.paragraphs[0]
    p.text = "Hardware Setup"
    p.font.size = Pt(18)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    vtb = _tb(sl, video_left, video_top - Inches(0.45), video_w, Inches(0.35))
    p = vtb.text_frame.paragraphs[0]
    p.text = "Baseline Trim-Transition Demo"
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    _card(sl, photo_card_left, photo_card_top, photo_card_w, photo_card_h, WHITE)
    _add_picture_fit_center(
        sl,
        photo_png,
        photo_card_left,
        photo_card_top,
        photo_card_w,
        photo_card_h,
    )

    _video(
        sl,
        video_path,
        video_left,
        video_top,
        video_w,
        video_h,
        poster_path=poster_png or photo_png,
    )
finally:
    if photo_png and os.path.exists(photo_png):
        os.unlink(photo_png)
    if poster_png and os.path.exists(poster_png):
        os.unlink(poster_png)

_callout(
    sl,
    Inches(2.0),
    Inches(8.95),
    Inches(16.0),
    Inches(1.0),
    "Baseline scenario: X-Plane pitch 0\u00b0 \u2192 +5\u00b0 \u2192 hold \u2192 return.  "
    "Platform tilts accordingly; PX4 senses in real time.",
    20,
)

_sn(sl, 4)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — System Architecture Overview
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "How the Physical Closed Loop Works")

fig_left = Inches(2.7)
fig_top = Inches(3.4)
_architecture_figure(sl, fig_left, fig_top, height=Cm(13.0))

callout_x = Inches(12.2)
callout_top = Inches(3.4)
steps = [
    ("1", "X-Plane publishes simulated\nvehicle state over UDP."),
    ("2", "Host converts that state into\nplatform pose commands."),
    ("3", "The Stewart platform creates\nreal inertial motion for PX4."),
    ("4", "PX4 attitude returns to host;\nhost PID injects control to X-Plane."),
]
for i, (num, en) in enumerate(steps):
    y = callout_top + Inches(i * 1.3)

    ntb = _tb(sl, callout_x, y, Inches(0.6), Inches(0.5))
    p = ntb.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(30)
    p.font.color.rgb = BLUE
    p.font.bold = True

    tb = _tb(sl, callout_x + Inches(0.7), y+Cm(0.1), Inches(6.0), Inches(1.3))
    ttf = tb.text_frame
    ttf.word_wrap = True
    p = ttf.paragraphs[0]
    p.text = en
    p.font.size = Pt(24)
    p.font.color.rgb = DARK
    p.line_spacing = Pt(24)

ftb = _tb(sl, Inches(2), Inches(9.3), Inches(16), Inches(0.6))
p = ftb.text_frame.paragraphs[0]
p.text = "전체 루프 구조를 먼저 확인하고, 이어지는 네 슬라이드에서 블록별 역할과 데이터 흐름을 살펴봅니다."
p.font.size = Pt(20)
p.font.color.rgb = BLUE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

_sn(sl, 5)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — Block Detail: X-Plane
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl, tf6 = _block_slide_base(prs, blank, 6, "Block 1 — X-Plane Simulator", 1)

p = tf6.paragraphs[0]
p.text = "\nROLE"
p.font.size = Pt(22)
p.font.color.rgb = BLUE
p.font.bold = True
_bullet(tf6, "기체 동역학 연산 및 반복 가능한 시나리오 환경 제공", None, 20, Pt(8))
_bullet(tf6, "UDP로 기체 상태(자세·위치·속도)를 호스트에 송신", None, 20, Pt(8))

_p(tf6, "", 8, sp=Pt(16))
_p(tf6, "KEY POINT", 22, BLUE, True, Pt(8))
_bullet(tf6, "기존에는 X-Plane이 자체적으로 루프를 닫았음", None, 20, Pt(8))
_bullet(tf6, "이 구조에서는 PX4 자세가 호스트 PID를 거쳐\n  제어 입력으로 되돌아옴", None, 20, Pt(8))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — Block Detail: Host
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl, tf7 = _block_slide_base(prs, blank, 7, "Block 2 — Python Host / Bridge", 2)

p = tf7.paragraphs[0]
p.text = "\nRECEIVE"
p.font.size = Pt(22)
p.font.color.rgb = BLUE
p.font.bold = True
_bullet(tf7, "X-Plane 상태 (UDP)", None, 20, Pt(8))
_bullet(tf7, "PX4 자세 추정값 (MAVLink / USB)", None, 20, Pt(8))

_p(tf7, "", 8, sp=Pt(16))
_p(tf7, "PROCESS", 22, BLUE, True, Pt(8))
_bullet(tf7, "시뮬레이터 상태 → 플랫폼 pose 명령으로 변환", None, 20, Pt(8))
_bullet(tf7, "모든 스트림에 공통 타임스탬프를 부여하여 로깅", None, 20, Pt(8))

_p(tf7, "", 8, sp=Pt(16))
_p(tf7, "RETURN", 22, BLUE, True, Pt(8))
_bullet(tf7, "pose 명령 → 플랫폼 (시리얼)", None, 20, Pt(8))
_bullet(tf7, "PX4 자세 → Host PID → X-Plane 제어 입력 주입 (UDP)", None, 20, Pt(8))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — Block Detail: Stewart Platform
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl, tf8 = _block_slide_base(prs, blank, 8, "Block 3 — Stewart Platform Actuation", 3)

p = tf8.paragraphs[0]
p.text = "\nROLE"
p.font.size = Pt(22)
p.font.color.rgb = BLUE
p.font.bold = True
_bullet(tf8, "호스트에서 시리얼로 수신한 pose 명령을 물리적 모션으로 변환", None, 20, Pt(8))
_bullet(tf8, "가동범위 포화 한계를 적용 (saturator)", None, 20, Pt(8))

_p(tf8, "", 8, sp=Pt(16))
_p(tf8, "RETURN", 22, BLUE, True, Pt(8))
_bullet(tf8, "명령 수신 ACK + 처리 소요 시간", None, 20, Pt(8))
_bullet(tf8, "포화 플래그 및 스케일링 계수 α", None, 20, Pt(8))

_p(tf8, "", 8, sp=Pt(16))
_p(tf8, "KEY POINT", 22, BLUE, True, Pt(8))
_bullet(tf8, "저비용 50 Hz RC 서보 | Serial RTT ~22 ms", None, 20, Pt(8))
_bullet(tf8, "이 블록 이후부터 PX4가 경험하는 것은 실제 관성 운동", None, 20, Pt(8))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 — Block Detail: PX4
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl, tf9 = _block_slide_base(prs, blank, 9, "Block 4 — PX4 Sensing and Return", 4)

p = tf9.paragraphs[0]
p.text = "\nSENSING"
p.font.size = Pt(22)
p.font.color.rgb = BLUE
p.font.bold = True
_bullet(tf9, "온보드 IMU가 플랫폼의 실제 운동을 직접 감지", None, 20, Pt(8))
_bullet(tf9, "PX4 추정기가 실시간으로 자세를 산출", None, 20, Pt(8))

_p(tf9, "", 8, sp=Pt(16))
_p(tf9, "RETURN", 22, BLUE, True, Pt(8))
_bullet(tf9, "추정된 자세를 MAVLink/USB로 호스트에 전송", None, 20, Pt(8))
_bullet(tf9, "호스트 PID가 조종면 명령 계산 후 X-Plane에 주입", None, 20, Pt(8))

_p(tf9, "", 8, sp=Pt(16))
_p(tf9, "KEY POINT", 22, BLUE, True, Pt(8))
_bullet(tf9, "유일한 물리적 관측 주체 — 추정기 반응 자체가 검증 목표", None, 20, Pt(8))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 — Experiment Design
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "What We Tested and How We Measured It")

col_w = Inches(7.5)
zone_top = Inches(3.4)
lx = Inches(1.8)
rx = Inches(10.7)

ltb = _tb(sl, lx, zone_top, col_w, Inches(6.5))
ltf = ltb.text_frame
ltf.word_wrap = True
p = ltf.paragraphs[0]
p.text = "Baseline Scenario"
p.font.size = Pt(26)
p.font.color.rgb = BLUE
p.font.bold = True
_p(ltf, "Pitch trim transition:", 22, DARK, True, Pt(18))
_p(ltf, "0\u00b0 \u2192 +5\u00b0 \u2192 hold \u2192 0\u00b0", 26, BLUE, True, Pt(6))
_p(ltf, "", 8, sp=Pt(8))
_bullet(ltf, "항상 동일 Honolulu (PHNL) 10nm Approach", "", 20, Pt(8))
_p(ltf, "", 8, sp=Inches(0.09))
_bullet(ltf, "Stewart 플랫폼의 포화 플래그(SAT) 감시;", "", 20, Pt(8))
_p(ltf, "", 8, sp=Inches(0.09))
_bullet(ltf, "조건별 9회 반복 실행으로 통계적 유의미성 확보", "", 20, Pt(8))
_p(ltf, "", 8, sp=Inches(0.09))
_bullet(ltf, "z=+20 mm에서 중립자세 설정 (최대 회전)", "", 20, Pt(8))

div = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.85), Inches(3.6), Pt(2), Inches(5.4))
div.fill.solid()
div.fill.fore_color.rgb = BLUE
div.line.fill.background()

rtb = _tb(sl, rx, zone_top, col_w, Inches(6.5))
rtf = rtb.text_frame
rtf.word_wrap = True
p = rtf.paragraphs[0]
p.text = "Logged Signals"
p.font.size = Pt(26)
p.font.color.rgb = BLUE
p.font.bold = True
for en in [
    "X-Plane attitude",
    "PX4 자세 추정값",
    "플랫폼 pose command",
    "ACK timing + host timestamps",
]:
    _bullet(rtf, en, None, 20, Pt(10))

_p(rtf, "", 8, sp=Pt(14))
_p(rtf, "Extracted Metrics", 24, BLUE, True, Pt(6))
for en in [
    "Tracking error (raw)",
    "Bias-corrected tracking error",
    "Sample age",
    "Serial RTT / end-to-end delay",
]:
    _bullet(rtf, en, None, 20, Pt(8))

_callout(sl, Inches(2.67), Inches(8.8), Inches(14.66), Inches(1.2),
         "시스템 자체 특성(지연·바이어스·장착 효과) 분리를 위해 clean baseline부터 수행", 20)

_sn(sl, 10)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11 — Result 1: Tracking
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "Result 1: Closed-Loop Tracking Is Achieved")

fig2 = os.path.join(FIG, "Fig2_PitchResponse.png")
_img(sl, fig2, Inches(1.0), Inches(3.4), width=Inches(11.5))

atb = _tb(sl, Inches(13.0), Inches(3.6), Inches(5.5), Inches(2.0))
atf = atb.text_frame
atf.word_wrap = True
p = atf.paragraphs[0]
p.text = ""

def _run(txt, bold=True, color=DARK):
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(20)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT_NAME
    return r

orange = RGBColor(0xE6, 0x7E, 0x22)
blue_dashed = RGBColor(0x4A, 0x90, 0xE2)

_run("점선", bold=True)
_run(": 스텝 입력. ")
_run("주황선", bold=True, color=orange)
_run(": FC pitch.\n")
_run("파란 점선", bold=True, color=blue_dashed)
_run(": X-Plane pitch.\n")
_run("파랑선", bold=True, color=BLUE)
_run(": +1.82° 보정된 X-Plane pitch이며 ")
_run("FC pitch", bold=True, color=orange)
_run("와 거의 일치함")

_p(atf, "", 8, sp=Pt(16))
_p(atf, "PX4는 일관된 위상 지연을 보이며 스텝을 추종하는 결과를 보임.", 20, DARK, sp=Pt(4))
_ko(atf, "", 20, Pt(8))

_p(atf, "", 8, sp=Pt(16))
_p(atf, "잔여 오프셋은 설명 가능한 바이어스임. 실험을 반복해도 일관되게 관찰됨.", 20, DARK, sp=Pt(4))
_ko(atf, "", 20, Pt(8))

_callout(sl, Inches(2.67), Inches(8.8), Inches(14.66), Inches(1.2),
         "플랫폼은 정량적 폐루프 검증이 가능할 정도로, 명령된 자세 전이를 충분히 정확하게 재현한다.", 20)

_sn(sl, 11)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12 — Result 2: Error Analysis
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "Result 2: Residual Error Is Interpretable")

pw = Inches(8.0)
fig3 = os.path.join(FIG, "Fig3_PitchError.png")
fig4 = os.path.join(FIG, "Fig4_MountingBias_TimeSeries.png")
_img(sl, fig3, Inches(1.0), Inches(3.4), width=pw)
_img(sl, fig4, Inches(10.0), Inches(3.4), width=pw)

cap1 = _tb(sl, Inches(1.0), Inches(7.6), pw, Inches(0.5))
p = cap1.text_frame.paragraphs[0]
p.text = "바이어스 보정을 적용하면 오차 분포가 0에 더 가까워진다."
p.font.size = Pt(20)
p.font.color.rgb = GREY
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

cap2 = _tb(sl, Inches(10.0), Inches(7.6), pw, Inches(0.5))
p = cap2.text_frame.paragraphs[0]
p.text = "단단히 고정한 실험에서는 장착 바이어스와 일치하는 안정적인 오프셋이 관찰된다."
p.font.size = Pt(20)
p.font.color.rgb = GREY
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

_callout(sl, Inches(2.67), Inches(8.8), Inches(14.66), Inches(1.2), "지배적인 잔여 오차는 무작위 불안정이 아니며, 고정 바이어스는 동적 추종 충실도와 분리해 해석해야 한다.", 20)

_sn(sl, 12)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13 — Result 3: Latency
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "Result 3: Latency Sets the Fidelity Limit")

fig5a = os.path.join(FIG, "Fig5a_Latency_TimeSeries.png")
fig5b = os.path.join(FIG, "Fig5b_Latency_Distribution_Baseline.png")
_img(sl, fig5a, Inches(1.0), Inches(3.4), width=Inches(8.5))
_img(sl, fig5b, Inches(10.0), Inches(3.4), width=Inches(8.5))

# cap5b = _tb(sl, Inches(10.0), Inches(7.95), Inches(8.5), Inches(0.45))
# p = cap5b.text_frame.paragraphs[0]
# p.text = "Fig5b: End-to-end latency distribution across baseline runs (#5-#9)"
# p.font.size = Pt(14)
# p.font.color.rgb = GREY
# p.font.italic = True
# p.alignment = PP_ALIGN.CENTER

chip_h = Inches(0.5)
chip_w = Inches(8.5)
chips = [
    ("X-Plane age: mean 27.2 / median 30.5 ms",  Inches(1.0),  Inches(7.35)),
    ("PX4 age: mean 18.9 / median 14.8 ms",      Inches(10.0), Inches(7.35)),
    ("Serial RTT: mean 21.7 / median 17.1 ms",    Inches(1.0),  Inches(7.85)),
    ("End-to-end: mean 48.9 / median 48.2 ms",    Inches(10.0), Inches(7.85)),
]
for text, x, y in chips:
    ctb = _tb(sl, x, y, chip_w, chip_h)
    p = ctb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

ctf = _callout(sl, Inches(2.67), Inches(8.8), Inches(14.66), Inches(1.2),"종단간 ~48 ms — 분포 폭이 크므로 유효 bandwidth는 전용 sine-sweep으로 정량화 예정", 20)

_sn(sl, 13)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14 — Takeaways
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "Takeaways")

col_w = Inches(7.5)
col_top = Inches(3.4)
lx = Inches(1.8)
rx = Inches(10.7)

ltb = _tb(sl, lx, col_top, col_w, Inches(5.5))
ltf = ltb.text_frame
ltf.word_wrap = True
p = ltf.paragraphs[0]
p.text = "Key Results"
p.font.size = Pt(26)
p.font.color.rgb = BLUE
p.font.bold = True

takeaways = [
    ("통합: X-Plane + Stewart + PX4를\n하나의 물리적 폐루프 구조로 통합",
     "통합: 물리적 폐루프 HILS 아키텍처 구현"),
    ("충실도: RMSE ~0.46° 추종 오차,\nE2E 지연 ~51 ms로 정량 검증",
     "충실도: RMSE/E2E 수치로 정량 검증"),
    ("특성화: 플랫폼의 가능/불가능 범위를\n정량적으로 파악 \u2014 향후 확장의 전제 조건",
     "특성화: 플랫폼 성능 정량 파악 \u2014 확장의 전제"),
]
for i, (en, kr) in enumerate(takeaways):
    _p(ltf, f"{i + 1}.  {en}", 20, DARK, sp=Pt(14))

_p(ltf, "", 8, sp=Pt(16))
_p(ltf, "※  SW HILS·비행 시험의 대체가 아닌 보완적 검증 수단",
   18, RED, True, Pt(4))

div = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.85), Inches(3.5), Pt(2), Inches(5.0))
div.fill.solid()
div.fill.fore_color.rgb = BLUE
div.line.fill.background()

rtb = _tb(sl, rx, col_top, col_w, Inches(5.5))
rtf = rtb.text_frame
rtf.word_wrap = True
p = rtf.paragraphs[0]
p.text = "Next Steps"
p.font.size = Pt(26)
p.font.color.rgb = BLUE
p.font.bold = True

for en, kr in [
    ("시간 정렬후 오차 분석", "시간 정렬 기반 오차 분석"),
    ("반복성·히스테리시스 특성 분석",
     "반복성·히스테리시스 특성 분석"),
    ("작동영역 및 포화 한계 부근 스트레스 테스트",
     "작업 영역·포화 한계 스트레스 시험"),
    ("멀티로터 등 다른 기체로 확장하여 아키텍처 범용성 실증",
     "다기체 확장으로 범용성 실증"),
]:
    _p(rtf, f"\u2192  {en}", 20, DARK, sp=Pt(14))

_callout(sl, Inches(2.67), Inches(8.8), Inches(14.66), Inches(1.2),
         "SW HILS와 비행 시험 사이의 보완적 검증 수단 — 작업 영역·공력 하중 한계 인지 하에 물리적 관성 자극 제공", 20)

_sn(sl, 14)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 15 — End / Thank You
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)

tb = _tb(sl, Inches(2), Inches(3.0), Inches(16), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(56)
p.font.color.rgb = BLACK
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

tb2 = _tb(sl, Inches(2), Inches(5.0), Inches(16), Inches(1.0))
p = tb2.text_frame.paragraphs[0]
p.text = "감사합니다"
p.font.size = Pt(32)
p.font.color.rgb = KO_GREY
p.alignment = PP_ALIGN.CENTER

tb3 = _tb(sl, Inches(2), Inches(6.5), Inches(16), Inches(1.0))
p = tb3.text_frame.paragraphs[0]
p.text = "Questions welcome"
p.font.size = Pt(24)
p.font.color.rgb = BLUE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

_sn(sl, 15)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 16 — Backup: z=+20 mm Proof
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

sl = make_slide(prs, blank)
slide_title(sl, "Q&A")

rot_fig = os.path.join(FIG, "rot_worst_vs_z.png")
_img(sl, rot_fig, Inches(2.5), Inches(3.4), width=Inches(10.0))

ann = _tb(sl, Inches(13.5), Inches(3.6), Inches(5.5), Inches(4.0))
atf = ann.text_frame
atf.word_wrap = True
p = atf.paragraphs[0]
p.text = "z = +20 mm"
p.font.size = Pt(28)
p.font.color.rgb = BLUE
p.font.bold = True
_p(atf, "\u2192 available range at this z", 22, DARK, sp=Pt(8))
_p(atf, "worst-direction limit ~20\u00b0;\nsufficient margin for baseline trim test.", 20, DARK, sp=Pt(16))
_p(atf, "", 8, sp=Pt(24))
_p(atf, "z = 0 mm", 28, RED, True, Pt(8))
_p(atf, "\u2192 smaller available range", 22, DARK, sp=Pt(8))
_p(atf, "worst-direction limit ~10\u00b0;\nless margin than z=+20 mm.", 20, GREY, sp=Pt(16))

_sn(sl, 16)


## Slide 17 (Q&A cards) removed — no dedicated visuals; answers are in the script.


# ── Save ────────────────────────────────────────────────────────────

out_path = os.path.join(BASE, "KSAS_Spring2026_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
